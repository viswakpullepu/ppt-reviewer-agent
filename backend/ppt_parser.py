"""PowerPoint file parser for extracting content and metadata."""
import os
from pathlib import Path
from typing import List, Dict, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
import logging

logger = logging.getLogger(__name__)


class PowerPointParser:
    """Parse PowerPoint files and extract slides, text, and metadata."""

    def __init__(self, file_path: str):
        """Initialize parser with PowerPoint file.
        
        Args:
            file_path: Path to .pptx or .ppt file
        """
        self.file_path = file_path
        self.presentation = None
        self.slides_data = []
        self._load_presentation()

    def _load_presentation(self):
        """Load PowerPoint presentation from file."""
        try:
            if not os.path.exists(self.file_path):
                raise FileNotFoundError(f"File not found: {self.file_path}")
            
            self.presentation = Presentation(self.file_path)
            logger.info(f"Successfully loaded presentation: {self.file_path}")
        except Exception as e:
            logger.error(f"Error loading presentation: {str(e)}")
            raise

    def get_slide_count(self) -> int:
        """Get total number of slides."""
        return len(self.presentation.slides)

    def extract_all_slides(self) -> List[Dict]:
        """Extract content from all slides.
        
        Returns:
            List of dictionaries containing slide data
        """
        self.slides_data = []
        
        for idx, slide in enumerate(self.presentation.slides):
            slide_data = self._extract_slide_content(slide, idx + 1)
            self.slides_data.append(slide_data)
        
        logger.info(f"Extracted {len(self.slides_data)} slides")
        return self.slides_data

    def _extract_slide_content(self, slide, slide_number: int) -> Dict:
        """Extract content from a single slide.
        
        Args:
            slide: PowerPoint slide object
            slide_number: Slide number for reference
            
        Returns:
            Dictionary with slide content and metadata
        """
        slide_data = {
            "slide_number": slide_number,
            "text_content": "",
            "title": "",
            "subtitle": "",
            "bullets": [],
            "images_count": 0,
            "shapes_count": 0,
            "font_sizes": [],
            "word_count": 0,
            "layout_name": slide.slide_layout.name,
        }
        
        texts = []
        image_count = 0
        
        for shape in slide.shapes:
            # Extract text from shapes
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
                slide_data["text_content"] += shape.text + "\n"
                
                # Get font sizes
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size:
                                slide_data["font_sizes"].append(int(run.font.size.pt))
                
                # Identify title and subtitle
                if "title" in shape.name.lower():
                    slide_data["title"] = shape.text
                elif "subtitle" in shape.name.lower():
                    slide_data["subtitle"] = shape.text
                elif shape.name.startswith("List"):
                    # Extract bullet points
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slide_data["bullets"].append({
                                    "text": paragraph.text,
                                    "level": paragraph.level,
                                })
            
            # Count images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1
        
        slide_data["images_count"] = image_count
        slide_data["shapes_count"] = len(slide.shapes)
        slide_data["word_count"] = len(slide_data["text_content"].split())
        slide_data["text_content"] = slide_data["text_content"].strip()
        
        return slide_data

    def get_presentation_metadata(self) -> Dict:
        """Get presentation-level metadata.
        
        Returns:
            Dictionary with presentation metadata
        """
        core_props = self.presentation.core_properties
        
        return {
            "title": core_props.title or "Untitled",
            "author": core_props.author or "Unknown",
            "subject": core_props.subject or "",
            "keywords": core_props.keywords or "",
            "total_slides": self.get_slide_count(),
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
        }

    def get_text_summary(self) -> str:
        """Get concatenated text from all slides.
        
        Returns:
            Full text summary
        """
        if not self.slides_data:
            self.extract_all_slides()
        
        summary = ""
        for slide in self.slides_data:
            summary += f"Slide {slide['slide_number']}:\n"
            if slide["title"]:
                summary += f"Title: {slide['title']}\n"
            summary += slide["text_content"] + "\n\n"
        
        return summary

    def analyze_slide_density(self, slide_data: Dict) -> Dict:
        """Analyze content density of a slide.
        
        Args:
            slide_data: Slide data dictionary
            
        Returns:
            Dictionary with density analysis
        """
        word_count = slide_data["word_count"]
        bullet_count = len(slide_data["bullets"])
        image_count = slide_data["images_count"]
        
        # Simple scoring (0-100)
        text_density = min(word_count / 100, 1.0) * 100
        
        density_rating = "optimal"
        if word_count > 150:
            density_rating = "too_dense"
        elif word_count < 20 and image_count == 0:
            density_rating = "too_sparse"
        
        return {
            "word_count": word_count,
            "bullet_count": bullet_count,
            "image_count": image_count,
            "text_density_score": text_density,
            "density_rating": density_rating,
            "recommendation": self._get_density_recommendation(density_rating),
        }

    def _get_density_recommendation(self, rating: str) -> str:
        """Get recommendation based on density rating."""
        recommendations = {
            "optimal": "Slide has good text-to-visual balance",
            "too_dense": "Too much text. Consider condensing or splitting into multiple slides.",
            "too_sparse": "Too little content. Add more information or visuals.",
        }
        return recommendations.get(rating, "")

    def get_all_analysis(self) -> Dict:
        """Get comprehensive analysis of presentation.
        
        Returns:
            Complete analysis dictionary
        """
        if not self.slides_data:
            self.extract_all_slides()
        
        analysis = {
            "metadata": self.get_presentation_metadata(),
            "slides": [],
        }
        
        for slide in self.slides_data:
            slide_analysis = {
                **slide,
                "density_analysis": self.analyze_slide_density(slide),
            }
            analysis["slides"].append(slide_analysis)
        
        return analysis
